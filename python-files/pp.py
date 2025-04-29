from pptx import Presentation
from pptx.util import Inches
import os
import re

# المسار إلى المجلد الذي يحتوي على الصور
image_folder = "img"  # قم بتعديل هذا المسار

# دالة ترتيب رقمي طبيعي
def natural_key(text):
    return [int(s) if s.isdigit() else s.lower() for s in re.split(r'(\d+)', text)]

# إنشاء ملف بوربوينت جديد
prs = Presentation()

# أبعاد الشريحة (اختياري – يمكن تعديله حسب الحاجة)
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)


# جمع الصور وترتيبها رقميًا
image_files = sorted(
    [f for f in os.listdir(image_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif'))],
    key=natural_key
)

# إضافة كل صورة في شريحة جديدة
for image in image_files:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # شريحة فارغة
    image_path = os.path.join(image_folder, image)
    slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

# حفظ الملف
prs.save("output_presentation.pptx")

print("✅ تم إنشاء ملف بوربوينت بترتيب الصور الرقمي بنجاح.")