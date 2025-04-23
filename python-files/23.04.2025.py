import os
import re
import asyncio
import zipfile
import win32com.client
import pythoncom
from openpyxl import Workbook  # Import Workbook from openpyxl
import uuid
import logging
import img2pdf
from telegram import (
    Update,
    InlineKeyboardButton,
    InlineKeyboardMarkup,
    Bot,
    InputFile
)
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    ContextTypes,
    filters
)
from telegram.constants import ParseMode
from docx2pdf import convert
from PIL import Image, ImageDraw, ImageFont
from fpdf import FPDF
from docx import Document
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from PIL import Image
import tempfile
from io import BytesIO
import subprocess
import sys
import docx
from openpyxl import load_workbook
from pptx import Presentation
from pdf2docx import Converter
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import pytesseract
import PyPDF2
import platform
import shutil
import psutil
import qrcode
import cv2
import numpy as np
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from pptx.util import Inches
import io

# Tesseract yo‚Äòlini sozlash (Windows uchun)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Logging configuration
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('bot.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Bot token
TOKEN = os.getenv('TELEGRAM_BOT_TOKEN', '6537782031:AAGI3LLWlvxIji4AzNn2Y0bPMSgDrWag_6I')

# Constants
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB for DOCX to PDF
MAX_EXCEL_SIZE = 5 * 1024 * 1024  # 5MB
ZIP_BUTTON = "zip_button"
CLEAR_BUTTON = "clear_button"
OCR_BUTTON = "ocr_button"
COMPRESS_BUTTON = "compress_button"
COMPRESS_HIGH = "compress_high"
COMPRESS_MEDIUM = "compress_medium"
COMPRESS_MAX = "compress_max"
OCR_MESSAGE = "ocr_message"
OCR_TXT = "ocr_txt"
OCR_DOCX = "ocr_docx"

# Translation dictionaries
CYRILLIC_TO_LATIN = {
    '–∞': 'a', '–±': 'b', '–≤': 'v', '–≥': 'g', '–¥': 'd', '–µ': 'e', '—ë': 'yo',
    '–∂': 'j', '–∑': 'z', '–∏': 'i', '–π': 'y', '–∫': 'k', '–ª': 'l', '–º': 'm',
    '–Ω': 'n', '–æ': 'o', '–ø': 'p', '—Ä': 'r', '—Å': 's', '—Ç': 't', '—É': 'u',
    '—Ñ': 'f', '—Ö': 'x', '—Ü': 'ts', '—á': 'ch', '—à': 'sh', '—â': 'shch',
    '—ä': "'", '—ã': 'i', '—å': "'", '—ç': 'e', '—é': 'yu', '—è': 'ya',
    '—û': "o'", '“õ': 'q', '“ì': "g'", '“≥': 'h',
    '–ê': 'A', '–ë': 'B', '–í': 'V', '–ì': 'G', '–î': 'D', '–ï': 'E', '–Å': 'Yo',
    '–ñ': 'J', '–ó': 'Z', '–ò': 'I', '–ô': 'Y', '–ö': 'K', '–õ': 'L', '–ú': 'M',
    '–ù': 'N', '–û': 'O', '–ü': 'P', '–†': 'R', '–°': 'S', '–¢': 'T', '–£': 'U',
    '–§': 'F', '–•': 'X', '–¶': 'Ts', '–ß': 'Ch', '–®': 'Sh', '–©': 'Shch',
    '–™': "'", '–´': 'I', '–¨': "'", '–≠': 'E', '–Æ': 'Yu', '–Ø': 'Ya',
    '–é': "O'", '“ö': 'Q', '“í': "G'", '“≤': 'H'
}

LATIN_TO_CYRILLIC = {
    'a': '–∞', 'b': '–±', 'd': '–¥', 'e': '–µ', 'f': '—Ñ', 'g': '–≥',
    'h': '“≥', 'i': '–∏', 'j': '–∂', 'k': '–∫', 'l': '–ª', 'm': '–º',
    'n': '–Ω', 'o': '–æ', 'p': '–ø', 'q': '“õ', 'r': '—Ä', 's': '—Å',
    't': '—Ç', 'u': '—É', 'v': '–≤', 'x': '—Ö', 'y': '–π', 'z': '–∑',
    "'": '—ä', '`': '—ä', ' ª': '—ä', ' º': '—ä', '‚Äô': '—ä',
    'sh': '—à', 'ch': '—á', 'ye': '–µ', 'ya': '—è', 'yu': '—é',
    "g'": '“ì', "o'": '—û',
    'A': '–ê', 'B': '–ë', 'D': '–î', 'E': '–ï', 'F': '–§', 'G': '–ì',
    'H': '“≤', 'I': '–ò', 'J': '–ñ', 'K': '–ö', 'L': '–õ', 'M': '–ú',
    'N': '–ù', 'O': '–û', 'P': '–ü', 'Q': '“ö', 'R': '–†', 'S': '–°',
    'T': '–¢', 'U': '–£', 'V': '–í', 'X': '–•', 'Y': '–ô', 'Z': '–ó',
    'Sh': '–®', 'Ch': '–ß', 'Ye': '–ï', 'Ya': '–Ø', 'Yu': '–Æ',
    "G'": '“í', "O'": '–é', "O‚Äò": '–é',
    'SH': '–®', 'CH': '–ß'
}

# User data storage
user_data = {}

class UserData:
    def __init__(self):
        self.files = []
        self.status_message_id = None
        self.images = []
        self.waiting_for_images = False
        self.active_module = None
        self.lock = asyncio.Lock()
        self.ocr_text = None
        self.output_format = None
        self.compressed_sizes = {}  # Siqilgan fayl hajmlarini saqlash uchun

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    keyboard = [
        [InlineKeyboardButton("üìÑ DOCX üîÑ PDF", callback_data='docx_pdf'),
         InlineKeyboardButton("üìä Excel ‚û° PDF", callback_data='excel_pdf')],
        [InlineKeyboardButton("üé§ PPT ‚û° PDF", callback_data='ppt_pdf'),
         InlineKeyboardButton("üñº RASM(JPG)-> PDF", callback_data='jpg_pdf')],
        [InlineKeyboardButton("üî§ Kiril üîÑ Lotin", callback_data='translate_file'),
         InlineKeyboardButton("üî≤ QR Kod Generator", callback_data='qr_gen')],
        [InlineKeyboardButton("üì∑ QR Kod Skaner", callback_data='qr_scan'),
         InlineKeyboardButton("üíß PDF Watermark", callback_data='pdf_watermark')],
        [InlineKeyboardButton("üîí PDF Parol qo'yish", callback_data='pdf_protect'),
         InlineKeyboardButton("üîì PDF Parolni olib tashlash", callback_data='pdf_unprotect')],
        [InlineKeyboardButton("üìù Word ‚û° Excel", callback_data='word_excel'),
         InlineKeyboardButton("üóÇ Arxivlovchi", callback_data='file_zipper')],
        [InlineKeyboardButton("üìú OCR (Matn chiqarish)", callback_data='ocr'),
         InlineKeyboardButton("üóú Fayl siqish", callback_data='compress')],
        [InlineKeyboardButton("üë®‚Äçüíª Dasturchi haqida", callback_data='about')]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    welcome_text = '''
    *Assalomu alaykum!* üëã
    
    Quyidagi yangi funksiyalar mavjud:
    - üìä *Excel ‚û° PDF* - Excel jadvallarini PDFga aylantirish
    - üé§ *PPT ‚û° PDF* - PowerPoint fayllarini PDFga aylantirish
    - üî≤ *QR Kod Generator* - Matn yoki URLdan QR kod yaratish
    - üì∑ *QR Kod Skaner* - Rasm orqali QR koddan ma'lumot olish
    - üíß *PDF Watermark* - PDF fayllarga suv belgisi qo'shish
    - üîí *PDF Parol qo'yish* - PDF fayllarga himoya qo'yish
    - üîì *PDF Parolni olib tashlash* - Parolli PDFlarni ochish
    - üìù *Word ‚û° Excel* - Word hujjatidagi jadvallarni Excelga o'tkazish
    
    ‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!
    '''
    
    await update.message.reply_text(
        welcome_text,
        reply_markup=reply_markup,
        parse_mode="Markdown"
    )

def format_file_size(size_bytes: int) -> str:
    """Fayl hajmini odamlar o'qiy oladigan formatga aylantirish"""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024:
            return f"{size_bytes:.2f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.2f} TB"

def compress_image(input_path: str, output_path: str, quality: str) -> bool:
    """Rasm fayllarini siqib, yangi faylga saqlaydi"""
    try:
        quality_settings = {
            'high': 90,
            'medium': 70,
            'max': 50
        }
        
        with Image.open(input_path) as img:
            if img.mode in ('RGBA', 'LA', 'P'):
                img = img.convert('RGB')
            
            img.save(output_path, "JPEG", quality=quality_settings[quality], optimize=True)
        
        return os.path.exists(output_path)
    
    except Exception as e:
        logger.error(f"Rasmni siqishda xato: {str(e)}")
        return False

async def compress_file_temp(input_path: str, file_extension: str, quality: str) -> int:
    """Faylni vaqtincha siqib, hajmini hisoblash"""
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = os.path.join(temp_dir, f"temp_compressed{file_extension}")
            
            if file_extension == '.pdf':
                # PDF uchun
                pdf_doc = fitz.open(input_path)
                image_paths = []
                
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc.load_page(page_num)
                    zoom = {
                        'high': 300/72,
                        'medium': 150/72,
                        'max': 72/72
                    }[quality]
                    
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)
                    
                    img_path = os.path.join(temp_dir, f"temp_page_{page_num}.jpg")
                    pix.save(img_path, "jpeg", quality={
                        'high': 90,
                        'medium': 70,
                        'max': 50
                    }[quality])
                    
                    image_paths.append(img_path)
                
                pdf_doc.close()
                
                # Rasmlardan PDF yaratish
                with open(output_path, "wb") as f:
                    f.write(img2pdf.convert(image_paths))
                
                if not os.path.exists(output_path):
                    return 0
                
            else:
                # Rasm fayllari uchun
                with Image.open(input_path) as img:
                    if img.mode in ('RGBA', 'LA', 'P'):
                        img = img.convert('RGB')
                    
                    quality_val = {
                        'high': 90,
                        'medium': 70,
                        'max': 50
                    }[quality]
                    
                    img.save(output_path, "JPEG", quality=quality_val, optimize=True)
            
            return os.path.getsize(output_path)
    
    except Exception as e:
        logger.error(f"Vaqtincha siqishda xato: {str(e)}")
        return 0

async def button_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    query = update.callback_query
    await query.answer()
    
    user_id = query.from_user.id
    chat_id = query.message.chat_id
    
    async with asyncio.Lock():
        if user_id not in user_data:
            user_data[user_id] = UserData()
    
    if query.data == 'excel_pdf':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'excel_pdf'
        message = await query.edit_message_text(
            "üì§ *Excel faylingizni yuboring (.xls, .xlsx)*\n\n"
            "‚ùó Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Natija PDF formatida qaytariladi\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'ppt_pdf':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'ppt_pdf'
        message = await query.edit_message_text(
            "üì§ *PowerPoint faylingizni yuboring (.ppt, .pptx)*\n\n"
            "‚ùó Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Natija PDF formatida qaytariladi\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'qr_gen':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'qr_gen'
        message = await query.edit_message_text(
            "üî≤ *QR kod yaratish*\n\n"
            "QR kodga aylantirmoqchi bo'lgan matn yoki URL manzilini yuboring:\n\n"
            "Misol: https://google.com yoki Salom Dunyo!\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'qr_scan':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'qr_scan'
        message = await query.edit_message_text(
            "üì∑ *QR kod skanerlash*\n\n"
            "QR kodni o'qish uchun rasm yuboring:\n\n"
            "‚ùó Iltimos, quyidagilarga e'tibor bering:\n"
            "- Rasm aniq va yorqin bo'lishi kerak\n"
            "- QR kod rasmdagi asosiy obyekt bo'lishi kerak\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'pdf_watermark':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'pdf_watermark'
        message = await query.edit_message_text(
            "üíß *PDF Watermark qo'shish*\n\n"
            "Suv belgisi qo'shmoqchi bo'lgan PDF faylingizni yuboring:\n\n"
            "Keyin sizga suv belgisi matnini so'raymiz\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'pdf_protect':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'pdf_protect'
        message = await query.edit_message_text(
            "üîí *PDF faylga parol qo'yish*\n\n"
            "Parol qo'ymoqchi bo'lgan PDF faylingizni yuboring:\n\n"
            "Keyin sizga parolni so'raymiz\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'pdf_unprotect':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'pdf_unprotect'
        message = await query.edit_message_text(
            "üîì *PDF fayldan parolni olib tashlash*\n\n"
            "Parolini olib tashlamoqchi bo'lgan PDF faylingizni yuboring:\n\n"
            "Keyin sizga parolni so'raymiz (agar ma'lum bo'lsa)\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'word_excel':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'word_excel'
        message = await query.edit_message_text(
            "üìù *Word dan Excelga o'tkazish*\n\n"
            "Excelga o'tkazmoqchi bo'lgan Word faylingizni yuboring:\n\n"
            "‚ùó Iltimos, quyidagilarga e'tibor bering:\n"
            "- Faylda jadvallar bo'lishi kerak\n"
            "- Fayl hajmi 10MB dan oshmasin\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'compress':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'compress'
            user_data[user_id].compressed_sizes = {}  # Hajmlarni tozalash
            user_data[user_id].temp_file_path = None  # Fayl yo‚Äòlini tozalash
        message = await query.edit_message_text(
            "üì§ *JPG, PNG yoki PDF faylingizni yuboring*\n\n"
            "‚ùó Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Fayl siqilgandan so‚Äòng sifat tanlash imkoniyati beriladi\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug‚Äòiga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data in [COMPRESS_HIGH, COMPRESS_MEDIUM, COMPRESS_MAX]:
        compression_level = {
            COMPRESS_HIGH: 'high',
            COMPRESS_MEDIUM: 'medium',
            COMPRESS_MAX: 'max'
        }[query.data]
    
        await query.answer("Fayl siqilmoqda...")
    
        try:
            # Oldingi xabarni o'chirish
            await context.bot.delete_message(chat_id=chat_id, message_id=query.message.message_id)
        
            # Fayl ma'lumotlarini olish
            async with user_data[user_id].lock:
                if not user_data[user_id].temp_file_path or not os.path.exists(user_data[user_id].temp_file_path):
                    await context.bot.send_message(chat_id, "‚ùå Avval fayl yuboring yoki fayl topilmadi!")
                    return
            
                input_path = user_data[user_id].temp_file_path
                file_name = user_data[user_id].temp_file_name
                file_extension = os.path.splitext(file_name)[1].lower()
                temp_dir = user_data[user_id].temp_dir
        
            processing_msg = await context.bot.send_message(chat_id, "‚è≥ Fayl siqilmoqda, iltimos kuting...")
        
            with tempfile.TemporaryDirectory() as output_temp_dir:
                output_path = os.path.join(output_temp_dir, f"compressed_{file_name}")
            
                if file_extension == '.pdf':
                    success = compress_pdf(input_path, output_path, compression_level)
                else:
                    success = compress_image(input_path, output_path, compression_level)
                if not success:
                    raise Exception("Faylni siqish muvaffaqiyatsiz tugadi")
                                    
                    
                # Natijani yuborish
                with open(output_path, 'rb') as result_file:
                    await context.bot.send_document(
                        chat_id=chat_id,
                        document=result_file,
                        filename=f"compressed_{file_name}",
                        caption=f"‚úÖ Fayl siqildi ({compression_level})!\nüåê @Convert_filesbot"
                    )
        
            await context.bot.delete_message(chat_id, processing_msg.message_id)
    
        except Exception as e:
            logger.error(f"Faylni siqishda xato: {str(e)}", exc_info=True)
            await context.bot.send_message(
                chat_id=chat_id,
                text=f"‚ùå Faylni siqishda xatolik yuz berdi: {str(e)}"
            )
    
        finally:
            # Tozalash
            async with user_data[user_id].lock:
                if 'input_path' in locals() and os.path.exists(input_path):
                    os.unlink(input_path)
                if 'temp_dir' in locals() and os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir, ignore_errors=True)
            
                user_data[user_id].temp_file_path = None
                user_data[user_id].temp_file_name = None
                user_data[user_id].compressed_sizes = {}
                user_data[user_id].temp_dir = None
        
            await asyncio.sleep(3)
            await return_to_main_menu(chat_id, context)
        

    # Qolgan button_handler qismlari o‚Äòzgarmagan holda saqlanadi
    elif query.data == 'ocr':
        async with user_data[user_id].lock:
            user_data[user_id].active_module = 'ocr'
        message = await query.edit_message_text(
            "üì§ *JPG, PNG yoki PDF faylingizni yuboring*\n\n"
            "‚ùó Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Matn aniq va o‚Äòqiladigan bo‚Äòlishi kerak\n"
            "- Natija formatini keyin tanlaysiz\n\n"
            "‚ÑπDasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = message.message_id
    
    elif query.data in [OCR_MESSAGE, OCR_TXT, OCR_DOCX]:
        async with user_data[user_id].lock:
            if not hasattr(user_data[user_id], 'ocr_text') or not user_data[user_id].ocr_text:
                await query.answer("Avval matn chiqarish uchun fayl yuboring!", show_alert=True)
                return
            output_format = query.data
            user_data[user_id].output_format = output_format
        
        try:
            await context.bot.delete_message(chat_id=chat_id, message_id=query.message.message_id)
        except Exception as e:
            logger.error(f"Format tanlash xabarini o‚Äòshirishda xato: {e}")
        
        await query.answer("Natija tayyorlanmoqda...")
        
        if output_format == OCR_MESSAGE:
            await context.bot.send_message(
                chat_id=chat_id,
                text=user_data[user_id].ocr_text[:4096],
                parse_mode="Markdown"
            )
        else:
            with tempfile.TemporaryDirectory() as temp_dir:
                if output_format == OCR_TXT:
                    output_path = os.path.join(temp_dir, "extracted_text.txt")
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(user_data[user_id].ocr_text)
                    filename = "extracted_text.txt"
                else:
                    output_path = os.path.join(temp_dir, "extracted_text.docx")
                    doc = Document()
                    doc.add_paragraph(user_data[user_id].ocr_text)
                    doc.save(output_path)
                    filename = "extracted_text.docx"
                
                with open(output_path, 'rb') as result_file:
                    await context.bot.send_document(
                        chat_id=chat_id,
                        document=result_file,
                        filename=filename,
                        caption=f"‚úÖ Matn muvaffaqiyatli chiqarildi!\nüåê @Convert_filesbot"
                    )
        
        async with user_data[user_id].lock:
            user_data[user_id].ocr_text = None
            user_data[user_id].output_format = None
        
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)
    
    elif query.data == 'docx_pdf':
        user_data[user_id].active_module = 'docx_pdf'
        message = await query.edit_message_text(
            "üì§ *PDF yoki DOCX faylingizni yuboring*\n\n"
            "‚ùó Iltimos, quyidagilarga e'tibor bering:\n"
            "- Fayl hajmi 10MB dan oshmasin\n"
            "- Faqat .docx yoki .pdf formatdagi fayllar qabul qilinadi\n"
            "- PDF fayllar DOCX ga, DOCX fayllar PDF ga aylantiriladi\n\n"
            "‚ÑπDasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            parse_mode="Markdown"
        )
        user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'jpg_pdf':
        user_data[user_id].active_module = 'jpg_pdf'
        user_data[user_id].images = []
        user_data[user_id].waiting_for_images = True
        message = await query.edit_message_text(
            "üì§PDF qilish kerak bo'lgan rasmlarni barchasini yuboring. "
            "Barcha rasmlarni yuborgach, 'PDF qilish' tugmasini bosing.\n\n"
            "‚ÑπDasturni qayta ishga tushurish uchun /start buyrug'iga bosing!"
        )
        user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'translate_file':
        user_data[user_id].active_module = 'kiril_lotin'
        message = await query.edit_message_text(
            "Iltimos, tarjima qilish uchun fayl yuboring:\n\n"
            "üìù Qo'llab-quvvatlanadigan formatlar:\n"
            "- Word hujjatlari (.doc, .docx)\n"
            "- Excel jadvallari (.xls, .xlsx)\n"
            "- PowerPoint prezentatsiyalari (.ppt, .pptx)\n"
            "- PDF fayllari (.pdf)\n\n"
            "Bot avtomatik tarjima yo'nalishini aniqlaydi!\n\n"
            "‚ÑπDasturni qayta ishga tushurish uchun /start buyrug'iga bosing!"
        )
        user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'file_zipper':
        user_data[user_id].active_module = 'file_zipper'
        user_data[user_id].files = []
        
        keyboard = [
            [InlineKeyboardButton("üóÇ Zip qilish (0)", callback_data=ZIP_BUTTON)]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        message = await query.edit_message_text(
            "üìÅ <b>Zip Botga xush kelibsiz!</b>\n\n"
            "üîπ Istalgan formatdagi fayllarni yuboring\n"
            "üîπ Barcha fayllar ro'yxatga qo'shiladi\n"
            "üîπ Pastdagi tugma orqali zip faylni oling\n\n"
            "üìÇ Joriy fayllar ro'yxati:\n"
            "Hozircha fayllar yo'q\n\n"
            "‚ÑπDasturni qayta ishga tushurish uchun /start buyrug'iga bosing!",
            reply_markup=reply_markup,
            parse_mode=ParseMode.HTML
        )
        user_data[user_id].status_message_id = message.message_id
    
    elif query.data == 'about':
        await query.edit_message_text(
            """
<b>üìù Telegram Bot</b>

<b>Versiya:</b> 2.0
<b>Ishlab chiqaruvchi:</b> Dilxush Bahodirov

<b>Funksiyalar:</b>
- DOCX ‚Üî PDF konvertor
- JPG dan PDF konvertor
- Kiril-Lotin avtomatik tarjimon
- Fayllarni zip arxiviga aylantirish

<b>Aloqa:</b> @Dilxush_Bahodirov
<b>‚ÑπDasturni qayta ishga tushurish uchun /start buyrug'iga bosing!</b>
            """,
            parse_mode=ParseMode.HTML
        )
    
    elif query.data == ZIP_BUTTON:
        await zip_files(update, context)
    
    elif query.data == CLEAR_BUTTON:
        await clear_files(update, context)
    
    elif query.data == 'create_pdf':
        await create_pdf(update, context)

async def handle_excel_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'excel_pdf':
        return
        
    try:
        if not update.message.document:
            await update.message.reply_text("‚ùå Iltimos, Excel fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.xlsx"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension not in ['.xls', '.xlsx']:
            await update.message.reply_text("‚ùå Noto'g'ri format! Iltimos, .xls yoki .xlsx fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"‚ùå Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("‚è≥ Excel fayl PDFga o'zgartirilmoqda...")
        
        # Vaqtinchalik papka yaratish
        temp_dir = tempfile.mkdtemp()
        try:
            # Fayl nomini tozalash
            import re
            safe_file_name = re.sub(r'[^\w\-\.]', '_', file_name)
            input_path = os.path.join(temp_dir, safe_file_name)
            await file.download_to_drive(input_path)
            
            # Fayl mavjudligini tekshirish
            if not os.path.exists(input_path):
                raise FileNotFoundError(f"Fayl topilmadi: {input_path}")
            
            # Faylni openpyxl bilan tekshirish (kontent mavjudligini aniqlash)
            try:
                wb_openpyxl = load_workbook(input_path, read_only=True)
                sheet_count = len(wb_openpyxl.sheetnames)
                if sheet_count == 0:
                    raise ValueError("Excel faylida varaqlar mavjud emas.")
                
                # Varaqlar kontentini tekshirish
                has_content = False
                for sheet in wb_openpyxl:
                    for row in sheet.iter_rows(max_row=10):  # Faqat birinchi 10 qatorni tekshirish
                        if any(cell.value for cell in row):
                            has_content = True
                            break
                    if has_content:
                        break
                wb_openpyxl.close()
                
                if not has_content:
                    raise ValueError("Excel faylida eksport qilinadigan kontent mavjud emas.")
                
                logger.info(f"openpyxl bilan tekshirildi: Varaqlar soni: {sheet_count}, Kontent mavjud: {has_content}")
            except Exception as openpyxl_error:
                logger.error(f"Faylni openpyxl bilan ochishda xato: {str(openpyxl_error)}, Fayl: {input_path}")
                raise Exception("Excel fayli shikastlangan yoki ochilmaydi.")
            
            output_path = os.path.join(temp_dir, f"{os.path.splitext(safe_file_name)[0]}.pdf")
            
            # Excel jarayonlarini oldindan tozalash
            for proc in psutil.process_iter(['name']):
                if proc.info['name'].lower() == 'excel.exe':
                    try:
                        proc.kill()
                        logger.info("Oldingi Excel jarayoni yopildi.")
                    except Exception as kill_error:
                        logger.warning(f"Excel jarayonini yopishda xato: {str(kill_error)}")
            
            pythoncom.CoInitialize()
            excel = None
            wb = None
            try:
                excel = win32com.client.Dispatch("Excel.Application")
                excel.Visible = False
                excel.DisplayAlerts = False
                
                # Faylni ochish
                try:
                    wb = excel.Workbooks.Open(input_path)
                except Exception as open_error:
                    logger.error(f"Excel faylini ochishda xato: {str(open_error)}, Fayl: {input_path}")
                    raise Exception("Excel faylini ochib bo‚Äòlmadi. Fayl shikastlangan yoki mos emas.")
                
                if wb is None:
                    raise ValueError("Excel fayli ochilmadi, wb ob'ekti None.")
                
                # Varaqlar sonini tekshirish
                sheet_count = wb.Worksheets.Count
                if sheet_count == 0:
                    raise ValueError("Excel faylida varaqlar mavjud emas.")
                
                logger.info(f"Fayl ochildi: {input_path}, Varaqlar soni: {sheet_count}")
                
                # Birinchi yondashuv: Varaqlarni PNG sifatida eksport qilish
                image_paths = []
                for i, ws in enumerate(wb.Worksheets, 1):
                    img_path = os.path.join(temp_dir, f"sheet_{i}.png")
                    try:
                        # Varaqni PNG sifatida eksport qilish
                        ws.ExportAsFixedFormat(2, img_path)
                        if os.path.exists(img_path) and os.path.getsize(img_path) > 0:
                            image_paths.append(img_path)
                            logger.info(f"Varaq {i} eksport qilindi: {img_path}")
                        else:
                            logger.warning(f"Varaq {i} eksport qilinmadi yoki bo‚Äòsh: {img_path}")
                    except Exception as export_error:
                        logger.error(f"Varaq {i} ni PNG eksport qilishda xato: {str(export_error)}")
                        continue
                
                # Agar PNG eksporti muvaffaqiyatsiz bo‚Äòlsa, to‚Äòg‚Äòridan-to‚Äòg‚Äòri PDF eksportini sinash
                if not image_paths:
                    logger.info("PNG eksporti muvaffaqiyatsiz. To‚Äòg‚Äòridan-to‚Äòg‚Äòri PDF eksporti sinab ko‚Äòrilmoqda.")
                    try:
                        wb.ExportAsFixedFormat(0, output_path)  # To‚Äòg‚Äòridan-to‚Äòg‚Äòri PDF sifatida
                        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                            logger.info(f"To‚Äòg‚Äòridan-to‚Äòg‚Äòri PDF eksport qilindi: {output_path}")
                        else:
                            raise Exception("To‚Äòg‚Äòridan-to‚Äòg‚Äòri PDF eksporti muvaffaqiyatsiz.")
                    except Exception as pdf_export_error:
                        logger.error(f"To‚Äòg‚Äòridan-to‚Äòg‚Äòri PDF eksportida xato: {str(pdf_export_error)}")
                        raise Exception("Hech qanday varaq eksport qilinmadi (PNG yoki PDF).")
                else:
                    # Rasmlarni PDF ga aylantirish
                    with open(output_path, "wb") as f:
                        f.write(img2pdf.convert(image_paths))
                    logger.info(f"PNG rasmlar PDF ga birlashtirildi: {output_path}")
                
                # PDF fayli yaratilganligini tekshirish
                if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                    raise Exception("PDF fayli yaratilmadi yoki bo‚Äòsh.")
                
            finally:
                # Excel resurslarini tozalash
                if wb is not None:
                    wb.Close(SaveChanges=False)
                if excel is not None:
                    excel.Quit()
                pythoncom.CoUninitialize()
                
                # Excel jarayonlarini qayta tekshirish va yopish
                for proc in psutil.process_iter(['name']):
                    if proc.info['name'].lower() == 'excel.exe':
                        try:
                            proc.kill()
                            logger.info("Excel jarayoni yopildi.")
                        except Exception as kill_error:
                            logger.warning(f"Excel jarayonini yopishda xato: {str(kill_error)}")
            
            # Natijani yuborish
            with open(output_path, 'rb') as result_file:
                await update.message.reply_document(
                    document=result_file,
                    caption=f"‚úÖ {os.path.splitext(safe_file_name)[0]}.pdf\n\nüåê @Convert_filesbot"
                )
            
            await context.bot.delete_message(chat_id, msg.message_id)
        
        finally:
            # Vaqtinchalik papkani tozalash
            shutil.rmtree(temp_dir, ignore_errors=True)
            logger.info(f"Vaqtinchalik papka tozalandi: {temp_dir}")
        
    except Exception as e:
        logger.error(f"Excel ni PDF ga o‚Äòtkazishda xato: {str(e)}, Fayl: {file_name}", exc_info=True)
        await update.message.reply_text(f"‚ùå Xatolik yuz berdi: {str(e)}")
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_ppt_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'ppt_pdf':
        return
        
    try:
        if not update.message.document:
            await update.message.reply_text("‚ùå Iltimos, PowerPoint fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "presentation.pptx"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension not in ['.ppt', '.pptx']:
            await update.message.reply_text("‚ùå Noto'g'ri format! Iltimos, .ppt yoki .pptx fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"‚ùå Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("‚è≥ PowerPoint fayl PDFga o'zgartirilmoqda...")
        
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, file_name)
            await file.download_to_drive(input_path)
            
            output_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.pdf")
            
            if file_extension == '.ppt':
                success = convert_ppt_to_pdf(input_path, output_path)
            else:
                success = convert_pptx_to_pdf(input_path, output_path)  # Keyingi bo‚Äòlimda yangilanadi
            
            if not success:
                raise Exception("PDF fayli yaratilmadi yoki bo'sh.")
            
            with open(output_path, 'rb') as result_file:
                await update.message.reply_document(
                    document=result_file,
                    filename=f"{os.path.splitext(file_name)[0]}.pdf",
                    caption=f"‚úÖ {os.path.splitext(file_name)[0]}.pdf\n\nüåê @Convert_filesbot"
                )
            
            await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"‚ùå Xatolik yuz berdi: {str(e)}")
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

def convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    try:
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        
        # Vaqtinchalik papka yaratish
        with tempfile.TemporaryDirectory() as temp_dir:
            # Slaydlarni rasm sifatida eksport qilish
            image_paths = []
            for i, slide in enumerate(presentation.Slides, 1):
                img_path = os.path.join(temp_dir, f"slide_{i}.png")
                slide.Export(img_path, "PNG", 1280, 720)
                image_paths.append(img_path)
            
            # Rasmlarni PDF ga aylantirish
            with open(pdf_path, "wb") as f:
                f.write(img2pdf.convert(image_paths))
            
        presentation.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()
        
        return os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0
    
    except Exception as e:
        logger.error(f"PPTX ni PDF ga aylantirishda xato: {str(e)}")
        return False

async def handle_qr_gen(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'qr_gen':
        return
        
    try:
        if not update.message.text:
            await update.message.reply_text("‚ùå Iltimos, QR kodga aylantirmoqchi bo'lgan matn yuboring!")
            return
            
        text = update.message.text
        msg = await update.message.reply_text("‚è≥ QR kod yaratilmoqda...")
        
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=4,
        )
        qr.add_data(text)
        qr.make(fit=True)
        
        img = qr.make_image(fill_color="black", back_color="white")
        
        with BytesIO() as bio:
            img.save(bio, 'PNG')
            bio.seek(0)
            await update.message.reply_photo(
                photo=bio,
                caption=f"‚úÖ QR kod yaratildi!\nMatn: {text[:50]}{'...' if len(text) > 50 else ''}\nüåê @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"‚ùå QR kod yaratishda xatolik: {str(e)}")
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_qr_scan(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'qr_scan':
        return
        
    try:
        if not update.message.photo:
            await update.message.reply_text("‚ùå Iltimos, QR kod rasmini yuboring!")
            return
            
        photo = update.message.photo[-1]
        msg = await update.message.reply_text("‚è≥ QR kod skanerlash...")
        
        file = await photo.get_file()
        with BytesIO() as bio:
            await file.download_to_memory(out=bio)
            bio.seek(0)
            
            img = Image.open(bio)
            img_array = np.array(img)
            
            detector = cv2.QRCodeDetector()
            data, _, _ = detector.detectAndDecode(img_array)
            
            if data:
                await update.message.reply_text(
                    f"‚úÖ QR kod mazmuni:\n\n{data}\n\nüåê @Convert_filesbot",
                    disable_web_page_preview=True
                )
            else:
                await update.message.reply_text("‚ùå QR kod aniqlanmadi yoki o'qib bo'lmadi!")
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"‚ùå QR kodni o'qishda xatolik: {str(e)}")
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_pdf_watermark(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_watermark':
        return
        
    try:
        if not update.message.document:
            await update.message.reply_text("‚ùå Iltimos, PDF fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.pdf"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension != '.pdf':
            await update.message.reply_text("‚ùå Noto'g'ri format! Iltimos, PDF fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"‚ùå Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("‚è≥ PDF fayl yuklanmoqda...")
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = os.path.join(tempfile.gettempdir(), f"watermark_{user_id}.pdf")
            await file.download_to_drive(user_data[user_id].temp_file_path)
            user_data[user_id].active_module = 'pdf_watermark_text'
        
        await msg.edit_text("‚úÖ PDF fayl qabul qilindi!\n\nEndi suv belgisi uchun matn yuboring:")
        
    except Exception as e:
        await update.message.reply_text(f"‚ùå Xatolik yuz berdi: {str(e)}")

async def handle_pdf_watermark_text(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_watermark_text':
        return
        
    try:
        if not update.message.text:
            await update.message.reply_text("‚ùå Iltimos, suv belgisi matnini yuboring!")
            return
            
        watermark_text = update.message.text
        msg = await update.message.reply_text("‚è≥ PDFga suv belgisi qo'shilmoqda...")
        
        input_path = user_data[user_id].temp_file_path
        output_path = os.path.join(tempfile.gettempdir(), f"watermarked_{user_id}.pdf")
        
        # Watermark qo'shish
        packet = BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.setFont("Helvetica", 50)
        can.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.3)
        can.rotate(45)
        can.drawString(100, 100, watermark_text)
        can.save()
        
        packet.seek(0)
        watermark_pdf = fitz.open("pdf", packet.getvalue())
        
        input_pdf = fitz.open(input_path)
        
        for page in input_pdf:
            page.show_pdf_page(page.rect, watermark_pdf, 0)
        
        input_pdf.save(output_path)
        input_pdf.close()
        watermark_pdf.close()
        
        with open(output_path, 'rb') as result_file:
            await update.message.reply_document(
                document=result_file,
                caption=f"‚úÖ Suv belgisi qo'shilgan PDF\nMatn: {watermark_text}\nüåê @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"‚ùå Suv belgisi qo'shishda xatolik: {str(e)}")
    finally:
        # Tozalash
        if 'input_path' in locals() and os.path.exists(input_path):
            os.unlink(input_path)
        if 'output_path' in locals() and os.path.exists(output_path):
            os.unlink(output_path)
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = None
        
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_pdf_protect(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_protect':
        return
        
    try:
        if not update.message.document:
            await update.message.reply_text("‚ùå Iltimos, PDF fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.pdf"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension != '.pdf':
            await update.message.reply_text("‚ùå Noto'g'ri format! Iltimos, PDF fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"‚ùå Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("‚è≥ PDF fayl yuklanmoqda...")
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = os.path.join(tempfile.gettempdir(), f"protect_{user_id}.pdf")
            await file.download_to_drive(user_data[user_id].temp_file_path)
            user_data[user_id].active_module = 'pdf_protect_password'
        
        await msg.edit_text("‚úÖ PDF fayl qabul qilindi!\n\nEndi parolni yuboring (6-32 belgi):")
        
    except Exception as e:
        await update.message.reply_text(f"‚ùå Xatolik yuz berdi: {str(e)}")

async def handle_pdf_protect_password(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_protect_password':
        return
        
    try:
        if not update.message.text:
            await update.message.reply_text("‚ùå Iltimos, parol yuboring!")
            return
            
        password = update.message.text
        if len(password) < 6 or len(password) > 32:
            await update.message.reply_text("‚ùå Parol 6-32 belgidan iborat bo'lishi kerak!")
            return
            
        msg = await update.message.reply_text("‚è≥ PDF faylga parol qo'yilmoqda...")
        
        input_path = user_data[user_id].temp_file_path
        output_path = os.path.join(tempfile.gettempdir(), f"protected_{user_id}.pdf")
        
        pdf_reader = PyPDF2.PdfReader(input_path)
        pdf_writer = PyPDF2.PdfWriter()
        
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
        
        pdf_writer.encrypt(password)
        
        with open(output_path, "wb") as f:
            pdf_writer.write(f)
        
        with open(output_path, 'rb') as result_file:
            await update.message.reply_document(
                document=result_file,
                caption=f"‚úÖ Parol bilan himoyalangan PDF\nüåê @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"‚ùå Parol qo'yishda xatolik: {str(e)}")
    finally:
        # Tozalash
        if 'input_path' in locals() and os.path.exists(input_path):
            os.unlink(input_path)
        if 'output_path' in locals() and os.path.exists(output_path):
            os.unlink(output_path)
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = None
        
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_pdf_unprotect(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_unprotect':
        return
        
    try:
        if not update.message.document:
            await update.message.reply_text("‚ùå Iltimos, parolli PDF fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.pdf"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension != '.pdf':
            await update.message.reply_text("‚ùå Noto'g'ri format! Iltimos, PDF fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"‚ùå Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("‚è≥ PDF fayl yuklanmoqda...")
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = os.path.join(tempfile.gettempdir(), f"unprotect_{user_id}.pdf")
            await file.download_to_drive(user_data[user_id].temp_file_path)
            user_data[user_id].active_module = 'pdf_unprotect_password'
        
        await msg.edit_text("‚úÖ PDF fayl qabul qilindi!\n\nAgar parol ma'lum bo'lsa, yuboring (agar parolni bilmasangiz, 'skip' yozing):")
        
    except Exception as e:
        await update.message.reply_text(f"‚ùå Xatolik yuz berdi: {str(e)}")

async def handle_pdf_unprotect_password(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'pdf_unprotect_password':
        return
        
    try:
        password = update.message.text if update.message.text.lower() != 'skip' else None
        
        msg = await update.message.reply_text("‚è≥ PDF fayldan parol olib tashlanmoqda...")
        
        input_path = user_data[user_id].temp_file_path
        output_path = os.path.join(tempfile.gettempdir(), f"unprotected_{user_id}.pdf")
        
        pdf_reader = PyPDF2.PdfReader(input_path)
        if pdf_reader.is_encrypted:
            if password:
                pdf_reader.decrypt(password)
            else:
                # Parolni bilmasa, oddiy usul bilan ochishga harakat qilish
                try:
                    pdf_reader.decrypt("")
                except:
                    raise Exception("Parol talab qilinadi va siz parolni yubormagansiz")
        
        pdf_writer = PyPDF2.PdfWriter()
        
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
        
        with open(output_path, "wb") as f:
            pdf_writer.write(f)
        
        with open(output_path, 'rb') as result_file:
            await update.message.reply_document(
                document=result_file,
                caption=f"‚úÖ Parolsiz PDF\nüåê @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await update.message.reply_text(f"‚ùå Parolni olib tashlashda xatolik: {str(e)}")
    finally:
        # Tozalash
        if 'input_path' in locals() and os.path.exists(input_path):
            os.unlink(input_path)
        if 'output_path' in locals() and os.path.exists(output_path):
            os.unlink(output_path)
        
        async with user_data[user_id].lock:
            user_data[user_id].temp_file_path = None
        
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_word_excel(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'word_excel':
        return
        
    try:
        if not update.message.document:
            await update.message.reply_text("‚ùå Iltimos, Word fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document.docx"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension not in ['.doc', '.docx']:
            await update.message.reply_text("‚ùå Noto'g'ri format! Iltimos, Word fayl yuboring.")
            return
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"‚ùå Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        msg = await update.message.reply_text("‚è≥ Word fayl Excelga o'tkazilmoqda...")
        
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, file_name)
            await file.download_to_drive(input_path)
            
            output_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.xlsx")
            
            # Convert .doc to .docx if necessary
            if file_extension == '.doc':
                temp_docx_path = os.path.join(temp_dir, "temp.docx")
                pythoncom.CoInitialize()
                word = win32com.client.Dispatch("Word.Application")
                doc = word.Documents.Open(input_path)
                doc.SaveAs(temp_docx_path, FileFormat=16)  # Save as DOCX
                doc.Close()
                word.Quit()
                pythoncom.CoUninitialize()
                input_path = temp_docx_path
            
            # Read the Word document
            doc = Document(input_path)
            
            # Create a new Excel workbook
            wb = Workbook()
            ws = wb.active
            ws.title = "Asosiy"
            
            # Track current row in Excel
            current_row = 1
            
            # Extract all content including tables
            for element in doc.element.body:
                if element.tag.endswith('tbl'):  # Jadval topildi
                    table = docx.table.Table(element, doc)
                    for row in table.rows:
                        for col_idx, cell in enumerate(row.cells):
                            # Jadval kataklarini Excelga yozish
                            ws.cell(row=current_row, column=col_idx+1).value = cell.text
                        current_row += 1
                    current_row += 1  # Jadval orasida bo'sh qator
                
                elif element.tag.endswith('p'):  # Oddiy paragraf
                    paragraph = docx.text.paragraph.Paragraph(element, doc)
                    if paragraph.text.strip():  # Faqat matn bor paragraflarni olish
                        ws.cell(row=current_row, column=1).value = paragraph.text
                        current_row += 1
            
            # Adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(cell.value)
                    except:
                        pass
                adjusted_width = (max_length + 2) * 1.2
                ws.column_dimensions[column_letter].width = adjusted_width
            
            # Save the Excel file
            wb.save(output_path)
            
            # Send the result
            with open(output_path, 'rb') as result_file:
                await update.message.reply_document(
                    document=result_file,
                    filename=f"{os.path.splitext(file_name)[0]}.xlsx",
                    caption=f"‚úÖ {os.path.splitext(file_name)[0]}.xlsx\n\nüåê @Convert_filesbot"
                )
            
            await context.bot.delete_message(chat_id, msg.message_id)
        
    except Exception as e:
        await msg.edit_text(f"‚ùå Xatolik yuz berdi: {str(e)}")
    finally:
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)

async def handle_compress(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    async with user_data[user_id].lock:
        if user_id not in user_data or user_data[user_id].active_module != 'compress':
            return
    
    try:
        # Fayl yuklash va tekshirish
        if update.message.document:
            document = update.message.document
            file_name = document.file_name or "document.pdf"
            file_extension = os.path.splitext(file_name)[1].lower()
            if file_extension != '.pdf':
                await update.message.reply_text("‚ùå Faqat PDF fayllar qabul qilinadi!")
                return
            file = await document.get_file()
        elif update.message.photo:
            file = await update.message.photo[-1].get_file()
            file_name = "photo.jpg"
            file_extension = '.jpg'
        else:
            await update.message.reply_text("‚ùå Iltimos, JPG yoki PDF fayl yuboring!")
            return
        
        if file.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"‚ùå Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        msg = await update.message.reply_text("‚è≥ Fayl tahlil qilinmoqda...‚ôªÔ∏è")
        
        # Faylni yuklash
        temp_dir = tempfile.mkdtemp()
        input_path = os.path.join(temp_dir, file_name)
        await file.download_to_drive(input_path)
        
        # PDF uchun maxsus tekshiruv
        if file_extension == '.pdf':
            try:
                with fitz.open(input_path) as doc:
                    if doc.is_encrypted:
                        await msg.edit_text("‚ùå PDF fayl parol bilan himoyalangan. Iltimos, parolsiz PDF yuboring.")
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return
            except Exception as e:
                await msg.edit_text(f"‚ùå PDF faylni ochib bo'lmadi: {str(e)}")
                shutil.rmtree(temp_dir, ignore_errors=True)
                return
        
        # Fayl hajmlarini hisoblash
        sizes = {}
        for quality in ['high', 'medium', 'max']:
            with tempfile.NamedTemporaryFile(suffix=file_extension, delete=False) as temp_file:
                output_path = temp_file.name
            
            if file_extension == '.pdf':
                success = compress_pdf(input_path, output_path, quality)
            else:
                success = compress_image(input_path, output_path, quality)
            
            if success:
                sizes[quality] = os.path.getsize(output_path)
                os.unlink(output_path)
            else:
                sizes[quality] = 0
        
        # Agar siqish muvaffaqiyatsiz bo'lsa
        if all(size == 0 for size in sizes.values()):
            raise Exception("Faylni siqish jarayonida xatolik yuz berdi")
        
        # Ma'lumotlarni saqlash
        async with user_data[user_id].lock:
            user_data[user_id].compressed_sizes = sizes
            user_data[user_id].temp_file_path = input_path
            user_data[user_id].temp_file_name = file_name
            user_data[user_id].temp_dir = temp_dir
        
        # Inline tugmalar
        keyboard = [
            [InlineKeyboardButton(f"Yuqori ({format_file_size(sizes['high'])})", callback_data=COMPRESS_HIGH)],
            [InlineKeyboardButton(f"O ªrta ({format_file_size(sizes['medium'])})", callback_data=COMPRESS_MEDIUM)],
            [InlineKeyboardButton(f"Maksimal ({format_file_size(sizes['max'])})", callback_data=COMPRESS_MAX)]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await msg.edit_text(
            "‚úÖ Fayl tahlil qilindi!\n"
            "Quyidagi siqish darajalaridan birini tanlang:\n"
            f"- Yuqori: {format_file_size(sizes['high'])}\n"
            f"- O ªrta: {format_file_size(sizes['medium'])}\n"
            f"- Maksimal: {format_file_size(sizes['max'])}\n\n"
            "‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug ªiga bosing!",
            reply_markup=reply_markup,
            parse_mode="Markdown"
        )
        
        async with user_data[user_id].lock:
            user_data[user_id].status_message_id = msg.message_id
    
    except Exception as e:
        logger.error(f"Faylni siqishda xato: {str(e)}", exc_info=True)
        error_msg = f"‚ùå Xatolik yuz berdi: {str(e)}"
        if "PDF" in str(e):
            error_msg += "\n\n‚ö†Ô∏è Iltimos, quyidagilarni tekshiring:"
            error_msg += "\n1. PDF fayl buzilmaganligiga ishonch hosil qiling"
            error_msg += "\n2. Fayl parol bilan himoyalanmaganligiga ishonch hosil qiling"
            error_msg += "\n3. Fayl hajmi 10MB dan oshmasligiga ahamiyat bering"
        
        if 'msg' in locals():
            await msg.edit_text(error_msg)
        else:
            await update.message.reply_text(error_msg)
        
        # Tozalash
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)
        async with user_data[user_id].lock:
            user_data[user_id].compressed_sizes = {}
            user_data[user_id].temp_file_path = None
            user_data[user_id].temp_file_name = None
            user_data[user_id].temp_dir = None

def compress_pdf(input_path: str, output_path: str, quality: str) -> bool:
    """PDF faylni siqib, yangi faylga saqlaydi"""
    try:
        # Siqish parametrlari
        quality_settings = {
            'high': {'dpi': 300, 'jpeg_quality': 90},
            'medium': {'dpi': 150, 'jpeg_quality': 70},
            'max': {'dpi': 72, 'jpeg_quality': 50}
        }
        settings = quality_settings[quality]

        # PDF faylni ochish
        pdf = fitz.open(input_path)
        image_paths = []

        # Har bir sahifani rasmga aylantirish
        for page in pdf:
            pix = page.get_pixmap(matrix=fitz.Matrix(settings['dpi']/72, settings['dpi']/72))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # Rasmni siqish
            img_buffer = BytesIO()
            img.save(img_buffer, format="JPEG", quality=settings['jpeg_quality'], optimize=True)
            
            # Vaqtinchalik faylga saqlash
            temp_img_path = os.path.join(tempfile.gettempdir(), f"temp_page_{len(image_paths)}.jpg")
            with open(temp_img_path, "wb") as f:
                f.write(img_buffer.getvalue())
            image_paths.append(temp_img_path)

        pdf.close()

        # Rasmlardan yangi PDF yaratish
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(image_paths))

        # Vaqtinchalik fayllarni tozalash
        for img_path in image_paths:
            if os.path.exists(img_path):
                os.unlink(img_path)

        return os.path.exists(output_path)

    except Exception as e:
        logger.error(f"PDF siqishda xato: {str(e)}")
        # Vaqtinchalik fayllarni tozalash
        if 'image_paths' in locals():
            for img_path in image_paths:
                if os.path.exists(img_path):
                    os.unlink(img_path)
        return False
    
async def handle_ocr(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    async with user_data[user_id].lock:
        if user_id not in user_data or user_data[user_id].active_module != 'ocr':
            return
    
    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
    
    try:
        if update.message.document:
            document = update.message.document
            file_name = document.file_name or "document"
            file_extension = os.path.splitext(file_name)[1].lower()
            if file_extension not in ['.pdf']:
                await update.message.reply_text("‚ùå Faqat PDF fayllar qabul qilinadi!")
                return
            file = await document.get_file()
        elif update.message.photo:
            file = await update.message.photo[-1].get_file()
            file_name = "photo.jpg"
            file_extension = '.jpg'
        else:
            await update.message.reply_text("‚ùå Iltimos, JPG, PNG yoki PDF fayl yuboring!")
            return
        
        if file.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"‚ùå Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        msg = await update.message.reply_text("‚è≥ Matn chiqarilmoqda...‚ôªÔ∏è")
        
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, file_name)
            await file.download_to_drive(input_path)
            
            extracted_text = ""
            if file_extension == '.pdf':
                pdf_doc = fitz.open(input_path)
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc.load_page(page_num)
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img_path = os.path.join(temp_dir, f"page_{page_num}.png")
                    pix.save(img_path)
                    extracted_text += pytesseract.image_to_string(Image.open(img_path), lang='uzb+eng')
                pdf_doc.close()
            else:
                extracted_text = pytesseract.image_to_string(Image.open(input_path), lang='uzb+eng')
            
            if not extracted_text.strip():
                await msg.edit_text("‚ùå Faylda matn topilmadi yoki o‚Äòqib bo‚Äòlmadi!")
                return
            
            async with user_data[user_id].lock:
                user_data[user_id].ocr_text = extracted_text
            
            keyboard = [
                [InlineKeyboardButton("Botga xabar sifatida", callback_data=OCR_MESSAGE)],
                [InlineKeyboardButton(".txt fayl sifatida", callback_data=OCR_TXT)],
                [InlineKeyboardButton(".docx fayl sifatida", callback_data=OCR_DOCX)]
            ]
            reply_markup = InlineKeyboardMarkup(keyboard)
            
            await msg.edit_text(
                "‚úÖ Matn muvaffaqiyatli chiqarildi!\n"
                "Qanday formatda olishni xohlaysiz?",
                reply_markup=reply_markup
            )
    
    except Exception as e:
        await msg.edit_text(f"‚ùå Xatolik yuz berdi: {str(e)}")

def convert_docx_to_pdf_alternative(docx_path, pdf_path):
    try:
        if sys.platform == "win32":
            libreoffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        else:
            libreoffice_path = "soffice"
        
        command = [
            libreoffice_path,
            '--headless',
            '--convert-to',
            'pdf',
            '--outdir',
            os.path.dirname(pdf_path),
            docx_path
        ]
        
        subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return True
    
    except (subprocess.CalledProcessError, FileNotFoundError):
        try:
            doc = Document(docx_path)
            pdf_doc = fitz.open()
            
            for para in doc.paragraphs:
                if para.text.strip():
                    page = pdf_doc.new_page()
                    page.insert_text((50, 50), para.text)
            
            pdf_doc.save(pdf_path)
            pdf_doc.close()
            return True
            
        except Exception as e:
            logger.error(f"Alternative DOCX to PDF conversion failed: {str(e)}")
            return False

async def handle_docx_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'docx_pdf':
        return
        
    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
    
    try:
        if not update.message.document:
            await update.message.reply_text("‚ùå Iltimos, DOCX yoki PDF fayl yuboring!")
            return
            
        document = update.message.document
        file_name = document.file_name or "document"
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if document.file_size > MAX_FILE_SIZE:
            await update.message.reply_text(
                f"‚ùå Fayl hajmi juda katta! Iltimos, {MAX_FILE_SIZE//(1024*1024)}MB dan kichik fayl yuboring."
            )
            return
        
        if file_extension not in ['.docx', '.pdf']:
            await update.message.reply_text(
                "‚ùå Noto'g'ri format! Iltimos, faqat .docx yoki .pdf formatdagi fayllarni yuboring."
            )
            return
        
        file = await update.message.document.get_file()
        
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, file_name)
            await file.download_to_drive(input_path)
            
            if file_extension == '.docx':
                output_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.pdf")
                msg = await update.message.reply_text("‚è≥ DOCX fayl PDF ga o'zgartirilmoqda...‚ôªÔ∏è")
                
                try:
                    convert(input_path, output_path)
                    
                    if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                        raise Exception("docx2pdf kutubxonasi orqali konvertatsiya muvaffaqiyatsiz")
                    
                except Exception as e:
                    logger.warning(f"docx2pdf failed, trying alternative: {str(e)}")
                    await msg.edit_text("‚ö†Ô∏è Asosiy konvertatsiya usuli muvaffaqiyatsiz, zaxira usul ishga tushirilmoqda...")
                    
                    success = convert_docx_to_pdf_alternative(input_path, output_path)
                    
                    if not success or not os.path.exists(output_path):
                        raise Exception("Barcha konvertatsiya usullari muvaffaqiyatsiz")
                
                with open(output_path, 'rb') as result_file:
                    await update.message.reply_document(
                        document=result_file,
                        caption=f"‚úÖ {os.path.splitext(file_name)[0]}.pdf\n\nüåê @Convert_filesbot"
                    )
                
                await context.bot.delete_message(chat_id, msg.message_id)
                
            elif file_extension == '.pdf':
                clean_name = re.sub(r'[\\/*?:"<>|]', "", file_name.replace(".pdf", ""))[:50]
                output_path = os.path.join(temp_dir, f"{clean_name}.docx")
                
                msg = await update.message.reply_text("‚è≥ PDF fayl DOCX ga o'zgartirilmoqda...‚ôªÔ∏è")
                
                try:
                    cv = Converter(input_path)
                    cv.convert(output_path, start=0, end=None, keep_layout=True, recognize=True)
                    cv.close()

                    if os.path.getsize(output_path) < 1024:
                        await msg.edit_text("‚ÑπÔ∏è Qo'shimcha ishlov berilmoqda...")
                        
                        pdf_doc = fitz.open(input_path)
                        doc = Document()
                        
                        first_page = pdf_doc.load_page(0)
                        section = doc.sections[0]
                        section.page_width = Inches(first_page.rect.width / 72)
                        section.page_height = Inches(first_page.rect.height / 72)
                        
                        for page_num in range(len(pdf_doc)):
                            page = pdf_doc.load_page(page_num)
                            blocks = page.get_text("dict")["blocks"]
                            
                            for block in blocks:
                                if "lines" in block:
                                    for line in block["lines"]:
                                        for span in line["spans"]:
                                            p = doc.add_paragraph()
                                            run = p.add_run(span["text"])
                                            font = run.font
                                            font.name = span["font"]
                                            font.size = Pt(span["size"])
                                
                                elif "image" in block:
                                    try:
                                        img_info = block["image"]
                                        image_stream = BytesIO(img_info["image"])
                                        doc.add_picture(image_stream, 
                                                      width=Inches(img_info["width"] / 72),
                                                      height=Inches(img_info["height"] / 72))
                                    except Exception:
                                        pass
                            
                            if page_num < len(pdf_doc) - 1:
                                doc.add_page_break()
                        
                        doc.save(output_path)
                        pdf_doc.close()

                    await context.bot.delete_message(chat_id, msg.message_id)
                    await update.message.reply_document(
                        document=open(output_path, 'rb'),
                        filename=f"{clean_name}.docx",
                        caption=f"‚úÖ {clean_name}.docx\n\nüåê @Convert_filesbot"
                    )
                
                except Exception as e:
                    await msg.edit_text(f"‚ùå Faylni DOCX ga o'zgartirishda xatolik: {str(e)}")
        
        if user_data[user_id].status_message_id:
            try:
                await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            except Exception as e:
                logger.error(f"Xabarni o'chirishda xato: {e}")
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)
    
    except Exception as e:
        await update.message.reply_text(f"‚ùå Xatolik yuz berdi: {str(e)}")

async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'jpg_pdf':
        return
    
    if not user_data[user_id].waiting_for_images:
        return
    
    if user_data[user_id].status_message_id:
        try:
            await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
        except Exception as e:
            logger.error(f"Xabarni o'chirishda xato: {e}")
    
    photo = update.message.photo[-1]
    
    try:
        photo_file = await photo.get_file()
        photo_bytes = BytesIO()
        await photo_file.download_to_memory(out=photo_bytes)
        photo_bytes.seek(0)
        
        with Image.open(photo_bytes) as img:
            img.verify()
        
        photo_bytes.seek(0)
        user_data[user_id].images.append(photo_bytes)
    except Exception as e:
        logger.error(f"Rasmni qayta ishlashda xato: {e}")
        await update.message.reply_text(f"Rasmni qayta ishlashda xato: {e}")
        return
    
    keyboard = [[InlineKeyboardButton("PDF qilish", callback_data='create_pdf')]]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    message = await update.message.reply_text(
        f"üñº Rasm qabul qilindi. Jami: {len(user_data[user_id].images)} ta rasm.\n"
        "Yana rasm yuborishingiz mumkin yoki PDF yaratish uchun tugmani üÜó bosing.",
        reply_markup=reply_markup
    )
    user_data[user_id].status_message_id = message.message_id

async def create_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    query = update.callback_query
    user_id = query.from_user.id
    chat_id = query.message.chat_id
    
    if user_id not in user_data or not user_data[user_id].images:
        await query.answer("Siz hech qanday rasm yubormagansiz!", show_alert=True)
        return
    
    await query.answer("PDF yaratilmoqda...‚ôªÔ∏è")
    
    try:
        if user_data[user_id].status_message_id:
            try:
                await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            except Exception as e:
                logger.error(f"Xabarni o'chirishda xato: {e}")
        
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        temp_files = []
        
        for img_bytes in user_data[user_id].images:
            try:
                img_bytes.seek(0)
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                    temp_path = temp_img.name
                    temp_files.append(temp_path)
                    
                    with Image.open(img_bytes) as img:
                        if img.mode in ('RGBA', 'LA', 'P'):
                            img = img.convert('RGB')
                        img.save(temp_path, "JPEG", quality=90)
                    
                    pdf.add_page()
                    pdf.image(temp_path, x=10, y=10, w=190)
            except Exception as e:
                logger.error(f"Rasmni PDFga qo'shishda xato: {e}")
                continue
        
        if pdf.pages:
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                pdf_path = temp_pdf.name
                pdf.output(pdf_path)
                
                with open(pdf_path, 'rb') as pdf_file:
                    await context.bot.send_document(
                        chat_id=chat_id,
                        document=pdf_file,
                        filename='images.pdf',
                        caption=f"Sizning {len(user_data[user_id].images)} ta rasmingizdan PDF fayl yaratildi‚úÖ \nüåê @Convert_filesbot"
                    )
            
            os.unlink(pdf_path)
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)
    except Exception as e:
        logger.error(f"PDF yaratishda xato: {e}")
        await context.bot.send_message(
            chat_id=chat_id,
            text=f"PDF yaratishda xatolik yuz berdi: {str(e)}"
        )
    finally:
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
            except Exception as e:
                logger.error(f"Vaqtinchalik faylni o'chirishda xato: {e}")
        
        if user_id in user_data:
            user_data[user_id].images = []
            user_data[user_id].waiting_for_images = False

async def handle_translation_doc(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.message.from_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'kiril_lotin':
        return
    
    progress_msg = None
    temp_dir = None
    output_path = None
    file_path = None

    try:
        if not update.message.document:
            await update.message.reply_text("‚ùå Iltimos, tarjima qilish uchun fayl yuboring!")
            return
        
        document = update.message.document
        file_name = document.file_name or "document"
        file_extension = os.path.splitext(file_name)[1].lower()

        if file_extension not in ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf']:
            await update.message.reply_text(
                "‚ö†Ô∏è Iltimos, quyidagi formatlardagi fayllarni yuboring:\n"
                "üìù Word (.doc, .docx)\n"
                "üìä Excel (.xls, .xlsx)\n"
                "üé§ PowerPoint (.ppt, .pptx)\n"
                "üìÑ PDF (.pdf)"
            )
            return

        progress_msg = await update.message.reply_text("‚è≥ Fayl yuklanmoqda...")

        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, file_name)
        
        file = await context.bot.get_file(document.file_id)
        await file.download_to_drive(file_path)
        
        original_extension = file_extension
        temp_converted_path = None
        
        if file_extension == '.pdf':
            await progress_msg.edit_text("‚è≥ PDF fayl DOCX ga o'tkazilmoqda...")
            temp_docx_path = os.path.join(temp_dir, f"temp_{os.path.splitext(file_name)[0]}.docx")
            
            try:
                cv = Converter(file_path)
                cv.convert(temp_docx_path, start=0, end=None, keep_layout=True, recognize=True)
                cv.close()
                
                if not os.path.exists(temp_docx_path) or os.path.getsize(temp_docx_path) < 1024:
                    await progress_msg.edit_text("‚ÑπÔ∏è Alternativ usul bilan qayta ishlanmoqda...")
                    pdf_doc = fitz.open(file_path)
                    doc = docx.Document()
                    
                    first_page = pdf_doc.load_page(0)
                    section = doc.sections[0]
                    section.page_width = Inches(first_page.rect.width / 72)
                    section.page_height = Inches(first_page.rect.height / 72)
                    
                    for page_num in range(len(pdf_doc)):
                        page = pdf_doc.load_page(page_num)
                        blocks = page.get_text("dict")["blocks"]
                        
                        for block in blocks:
                            if "lines" in block:
                                for line in block["lines"]:
                                    for span in line["spans"]:
                                        p = doc.add_paragraph()
                                        run = p.add_run(span["text"])
                                        font = run.font
                                        font.name = span["font"]
                                        font.size = Pt(span["size"])
                            
                            elif "image" in block:
                                try:
                                    img_info = block["image"]
                                    image_stream = BytesIO(img_info["image"])
                                    doc.add_picture(image_stream, 
                                                  width=Inches(img_info["width"] / 72),
                                                  height=Inches(img_info["height"] / 72))
                                except Exception:
                                    pass
                        
                        if page_num < len(pdf_doc) - 1:
                            doc.add_page_break()
                    
                    doc.save(temp_docx_path)
                    pdf_doc.close()
                
                file_path = temp_docx_path
                original_extension = '.docx'
                
            except Exception as e:
                logger.error(f"PDFni DOCXga o'tkazishda xato: {e}")
                raise Exception(f"PDF faylni DOCX formatiga o'tkazishda xatolik: {str(e)}")

        if file_extension in ['.doc', '.xls', '.ppt']:
            convert_functions = {
                '.doc': (convert_doc_to_docx, '.docx'),
                '.xls': (convert_xls_to_xlsx, '.xlsx'),
                '.ppt': (convert_ppt_to_pptx, '.pptx')
            }
            
            convert_func, new_ext = convert_functions[file_extension]
            await progress_msg.edit_text(f"‚è≥ {file_extension.upper()} fayl {new_ext.upper()} ga o'tkazilmoqda...")
            
            temp_converted_path = os.path.join(temp_dir, f"converted{new_ext}")
            if not convert_func(file_path, temp_converted_path):
                raise Exception(f"{file_extension.upper()} -> {new_ext.upper()} konvertatsiyada xatolik")
            
            file_path = temp_converted_path
            original_extension = new_ext

        await progress_msg.edit_text("‚è≥ Fayl tahlil qilinmoqda...")
        sample_text = extract_sample_text(file_path, original_extension)
        
        if not sample_text.strip():
            if original_extension in ['.xls', '.xlsx']:
                sample_text = "1"
            else:
                raise Exception("Faylda tarjima qilinadigan matn topilmadi")

        script = detect_script(sample_text)
        if script == 'cyrillic':
            translation_dict = CYRILLIC_TO_LATIN
            direction = 'cyr_lat'
            detected_script = "Kirill"
            target_script = "Lotin"
        else:
            translation_dict = LATIN_TO_CYRILLIC
            direction = 'lat_cyr'
            detected_script = "Lotin"
            target_script = "Kirill"

        output_filename = f"tarjima_{os.path.splitext(file_name)[0]}{original_extension}"
        output_path = os.path.join(temp_dir, output_filename)
        
        await progress_msg.edit_text(f"‚è≥ {detected_script} ‚û°Ô∏è {target_script} tarjima qilinmoqda...")
        
        if not translate_file(file_path, output_path, translation_dict, direction, original_extension):
            raise Exception("Tarjima jarayonida xatolik yuz berdi")

        if file_extension == '.pdf':
            pdf_output_path = os.path.join(temp_dir, f"tarjima_{os.path.splitext(file_name)[0]}.pdf")
            await progress_msg.edit_text("‚è≥ Tarjima qilingan DOCX PDF ga o'tkazilmoqda...")
            
            try:
                convert(output_path, pdf_output_path)
                output_path = pdf_output_path
                output_filename = f"tarjima_{os.path.splitext(file_name)[0]}.pdf"
            except Exception as e:
                logger.error(f"DOCXni PDFga o'tkazishda xato: {e}")
                raise Exception(f"Tarjima qilingan DOCXni PDFga o'tkazishda xatolik: {str(e)}")

        await progress_msg.edit_text("‚úÖ Tarjima tugallandi! Fayl yuborilmoqda...")
        
        with open(output_path, 'rb') as result_file:
            await update.message.reply_document(
                document=result_file,
                filename=output_filename,
                caption=f"‚úÖ {detected_script} ‚û°Ô∏è {target_script} tarjima qilindi\n"
                        f"üìÑ Original fayl: {file_name}\n"
                        f"üåê @Convert_filesbot"
            )
        
        await context.bot.delete_message(chat_id, progress_msg.message_id)
        
    except Exception as e:
        logger.error(f"Faylni qayta ishlashda xatolik: {e}", exc_info=True)
        
        error_msg = "‚ùå Faylni qayta ishlashda xatolik yuz berdi.\n"
        
        if "The file is corrupt" in str(e):
            error_msg += "Fayl buzilgan yoki parol bilan himoyalangan.\n"
        elif "The password is incorrect" in str(e):
            error_msg += "Fayl parol bilan himoyalangan.\n"
        elif "Faylda tarjima qilinadigan matn topilmadi" in str(e):
            error_msg += "Faylda tarjima qilinadigan matn topilmadi.\n"
        else:
            error_msg += f"Xato tafsilotlari: {str(e)}\n"
        
        error_msg += "\nIltimos, quyidagilarni tekshiring:\n"
        error_msg += "1. Fayl formati qo'llab-quvvatlanadimi?\n"
        error_msg += "2. Fayl buzilmaganligiga ishonch hosil qiling\n"
        error_msg += "3. Fayl paroli bilan himoyalanmaganligiga ishonch hosil qiling\n\n"
        error_msg += "Agar muammo davom etsa, /start buyrug'i orqali qayta urinib ko'ring."
        
        if progress_msg:
            try:
                await progress_msg.edit_text(error_msg)
            except:
                await update.message.reply_text(error_msg)
        else:
            await update.message.reply_text(error_msg)
        
    finally:
        try:
            if output_path and os.path.exists(output_path):
                os.unlink(output_path)
            if file_path and os.path.exists(file_path):
                os.unlink(file_path)
            if temp_dir and os.path.exists(temp_dir):
                for filename in os.listdir(temp_dir):
                    file_path = os.path.join(temp_dir, filename)
                    try:
                        if os.path.isfile(file_path):
                            os.unlink(file_path)
                    except Exception as e:
                        logger.error(f"{file_path} ni o'chirishda xato: {e}")
                os.rmdir(temp_dir)
        except Exception as e:
            logger.error(f"Tozalashda xato: {e}")

    await asyncio.sleep(3)
    await return_to_main_menu(chat_id, context)
            
async def handle_zip_file(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    chat_id = update.message.chat_id
    
    if user_id not in user_data or user_data[user_id].active_module != 'file_zipper':
        return
    
    try:
        if update.message.document:
            file = await update.message.document.get_file()
            original_name = update.message.document.file_name or f"file_{len(user_data[user_id].files)}"
        elif update.message.photo:
            photo = update.message.photo[-1]
            file = await photo.get_file()
            original_name = f"üñº Rasm_{len(user_data[user_id].files)}.jpg"
        elif update.message.video:
            file = await update.message.video.get_file()
            original_name = update.message.video.file_name or f"üé¨ Video_{len(user_data[user_id].files)}.mp4"
        elif update.message.audio:
            file = await update.message.audio.get_file()
            original_name = update.message.audio.file_name or f"üéµ Audio_{len(user_data[user_id].files)}.mp3"
        else:
            await update.message.reply_text("‚ùå Qo'llab-quvvatlanmaydigan fayl turi!")
            return
        
        file_path = f"temp_{user_id}_{len(user_data[user_id].files)}"
        await file.download_to_drive(file_path)
        
        unique_name = generate_unique_name(original_name)
        user_data[user_id].files.append((file_path, original_name, unique_name))
        
        files_list = "\n".join([f"üìÑ {i+1}. {name}" for i, (_, name, _) in enumerate(user_data[user_id].files)])
        
        keyboard = [
            [InlineKeyboardButton(f"üóÇ Zip qilish ({len(user_data[user_id].files)})", callback_data=ZIP_BUTTON)],
            [InlineKeyboardButton("üßπ Ro'yxatni tozalash", callback_data=CLEAR_BUTTON)]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        try:
            await context.bot.edit_message_text(
                chat_id=chat_id,
                message_id=user_data[user_id].status_message_id,
                text=f"üìÅ <b>Fayl qabul qilindi!</b>\n\n"
                     f"üìÇ Joriy fayllar ro'yxati: \n"
                     f"{files_list}\n\n"
                     f"‚ÑπÔ∏è Jami: {len(user_data[user_id].files)} ta fayl",
                reply_markup=reply_markup,
                parse_mode=ParseMode.HTML
            )
        except:
            message = await update.message.reply_text(
                f"üìÅ <b>Fayl qabul qilindi!</b>\n\n"
                f"üìÇ Joriy fayllar ro'yxati:\n"
                f"{files_list}\n\n"
                f"‚ÑπÔ∏è Jami: {len(user_data[user_id].files)} ta fayl",
                reply_markup=reply_markup,
                parse_mode=ParseMode.HTML
            )
            user_data[user_id].status_message_id = message.message_id
    
    except Exception as e:
        await update.message.reply_text(f"‚ùå Xatolik yuz berdi: {str(e)}")

async def zip_files(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    user_id = query.from_user.id
    chat_id = query.message.chat_id
    
    if user_id not in user_data or not user_data[user_id].files:
        await query.edit_message_text(text="‚ö†Ô∏è Ziplash uchun fayllarni yuboring!")
        return
    
    try:
        if user_data[user_id].status_message_id:
            try:
                await context.bot.delete_message(chat_id, user_data[user_id].status_message_id)
            except Exception as e:
                logger.error(f"Xabarni o'chirishda xato: {e}")
        
        zip_filename = f"Arxiv_{user_id}.zip"
        
        with zipfile.ZipFile(zip_filename, 'w') as zipf:
            for file_path, _, unique_name in user_data[user_id].files:
                zipf.write(file_path, unique_name)
        
        await context.bot.send_chat_action(chat_id=chat_id, action="upload_document")
        await context.bot.send_document(
            chat_id=chat_id,
            document=open(zip_filename, 'rb'),
            filename="archive.zip",
            caption=f"‚úÖ {len(user_data[user_id].files)} ta fayl arxivlandi!\n üåê @Convert_filesbot"
        )
        for file_path, _, _ in user_data[user_id].files:
            if os.path.exists(file_path):
                os.remove(file_path)
        user_data[user_id].files = []
        await asyncio.sleep(3)
        await return_to_main_menu(chat_id, context)
            
    except Exception as e:
        await query.edit_message_text(text=f"‚ùå Xatolik yuz berdi: {str(e)}")
    
    finally:
        if os.path.exists(zip_filename):
            os.remove(zip_filename)

async def clear_files(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    user_id = update.effective_user.id
    
    if user_id in user_data:
        for file_path, _, _ in user_data[user_id].files:
            if os.path.exists(file_path):
                os.remove(file_path)
        user_data[user_id].files = []
    
    keyboard = [
        [InlineKeyboardButton("üóÇ ARXIVLASH (0)", callback_data=ZIP_BUTTON)]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await query.edit_message_text(
        text="üìÅ <b>Ro'yxat tozalandi!</b>\n\n"
             "üîπ Yangi fayllar yuborishingiz mumkin\n"
             "üîπ Fayllar ro'yxati bo'sh\n\n"
             "‚ÑπÔ∏è Jami: 0 ta fayl",
        reply_markup=reply_markup,
        parse_mode=ParseMode.HTML
    )

async def return_to_main_menu(chat_id: int, context: ContextTypes.DEFAULT_TYPE, delay: int = 1):
    await asyncio.sleep(delay)
    keyboard = [
        [InlineKeyboardButton("üìÑ DOCX üîÑ PDF", callback_data='docx_pdf'),
         InlineKeyboardButton("üìä Excel ‚û° PDF", callback_data='excel_pdf')],
        [InlineKeyboardButton("üé§ PPT ‚û° PDF", callback_data='ppt_pdf'),
         InlineKeyboardButton("üñº RASM(JPG)-> PDF", callback_data='jpg_pdf')],
        [InlineKeyboardButton("üî§ Kiril üîÑ Lotin", callback_data='translate_file'),
         InlineKeyboardButton("üî≤ QR Kod Generator", callback_data='qr_gen')],
        [InlineKeyboardButton("üì∑ QR Kod Skaner", callback_data='qr_scan'),
         InlineKeyboardButton("üíß PDF Watermark", callback_data='pdf_watermark')],
        [InlineKeyboardButton("üîí PDF Parol qo'yish", callback_data='pdf_protect'),
         InlineKeyboardButton("üîì PDF Parolni olib tashlash", callback_data='pdf_unprotect')],
        [InlineKeyboardButton("üìù Word ‚û° Excel", callback_data='word_excel'),
         InlineKeyboardButton("üóÇ Arxivlovchi", callback_data='file_zipper')],
        [InlineKeyboardButton("üìú OCR (Matn chiqarish)", callback_data='ocr'),
         InlineKeyboardButton("üóú Fayl siqish", callback_data='compress')],
        [InlineKeyboardButton("üë®‚Äçüíª Dasturchi haqida", callback_data='about')]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    welcome_text = '''
    *Assalomu alaykum!* üëã
    
    Quyidagi yangi funksiyalar mavjud:
    - üìä *Excel ‚û° PDF* - Excel jadvallarini PDFga aylantirish
    - üé§ *PPT ‚û° PDF* - PowerPoint fayllarini PDFga aylantirish
    - üî≤ *QR Kod Generator* - Matn yoki URLdan QR kod yaratish
    - üì∑ *QR Kod Skaner* - Rasm orqali QR koddan ma'lumot olish
    - üíß *PDF Watermark* - PDF fayllarga suv belgisi qo'shish
    - üîí *PDF Parol qo'yish* - PDF fayllarga himoya qo'yish
    - üîì *PDF Parolni olib tashlash* - Parolli PDFlarni ochish
    - üìù *Word ‚û° Excel* - Word hujjatidagi jadvallarni Excelga o'tkazish
    
    ‚Ñπ Dasturni qayta ishga tushurish uchun /start buyrug'iga bosing!
    '''
    
    await context.bot.send_message(
        chat_id=chat_id,
        text=welcome_text,
        reply_markup=reply_markup,
        parse_mode="Markdown"
    )

def generate_unique_name(filename):
    name, ext = os.path.splitext(filename)
    return f"{name}_{uuid.uuid4().hex[:6]}{ext}"

def detect_script(text: str) -> str:
    cyrillic_chars = sum(1 for char in text if '\u0400' <= char <= '\u04FF')
    latin_chars = sum(1 for char in text if char.isalpha() and char.lower() in LATIN_TO_CYRILLIC)
    return 'cyrillic' if cyrillic_chars > latin_chars else 'latin'

def extract_sample_text(file_path: str, file_extension: str) -> str:
    sample_text = ""
    try:
        if file_extension == '.doc':
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            doc = word.Documents.Open(file_path)
            sample_text = doc.Content.Text[:1000]
            doc.Close()
            word.Quit()
            pythoncom.CoUninitialize()
            
        elif file_extension == '.docx':
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    sample_text += para.text + " "
            
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            sample_text += cell.text + " "
        
        elif file_extension == '.xls':
            pythoncom.CoInitialize()
            excel = win32com.client.Dispatch("Excel.Application")
            wb = excel.Workbooks.Open(file_path)
            
            for sheet in wb.Sheets:
                used_range = sheet.UsedRange
                for row in used_range.Rows:
                    for cell in row.Columns:
                        if cell.Value and str(cell.Value).strip():
                            sample_text += str(cell.Value) + " "
                            if len(sample_text) > 1000:
                                break
                if len(sample_text) > 1000:
                    break
            
            wb.Close()
            excel.Quit()
            pythoncom.CoUninitialize()
            
        elif file_extension == '.xlsx':
            wb = load_workbook(file_path, read_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell and str(cell).strip():
                            sample_text += str(cell) + " "
                            if len(sample_text) > 1000:
                                break
                if len(sample_text) > 1000:
                    break
        
        elif file_extension == '.ppt':
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(file_path)
            
            for slide in presentation.Slides:
                for shape in slide.Shapes:
                    if hasattr(shape, "TextFrame"):
                        if shape.TextFrame.HasText:
                            sample_text += shape.TextFrame.TextRange.Text + " "
                            if len(sample_text) > 1000:
                                break
                if len(sample_text) > 1000:
                    break
            
            presentation.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()
            
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        sample_text += shape.text + " "
                        if len(sample_text) > 1000:
                            break
                if len(sample_text) > 1000:
                    break
    
    except Exception as e:
        logger.error(f"Namuna matn olishda xatolik: {e}")
    
    return sample_text.strip()[:1000]

def translate_old_doc(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        doc = word.Documents.Open(input_path)
        
        for paragraph in doc.Paragraphs:
            paragraph.Range.Text = translate_text(paragraph.Range.Text, translation_dict, direction)
        
        new_output_path = output_path.replace('.doc', '.docx')
        doc.SaveAs2(new_output_path, FileFormat=16)
        doc.Close()
        word.Quit()
        pythoncom.CoUninitialize()
        
        if new_output_path != output_path:
            if os.path.exists(output_path):
                os.remove(output_path)
            os.rename(new_output_path, output_path)
        
        return True
    except Exception as e:
        logger.error(f".doc faylni tarjima qilishda xato: {e}")
        return False

def translate_old_xls(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        pythoncom.CoInitialize()
        excel = win32com.client.Dispatch("Excel.Application")
        wb = excel.Workbooks.Open(input_path)
        
        for sheet in wb.Sheets:
            used_range = sheet.UsedRange
            for row in used_range.Rows:
                for cell in row.Columns:
                    if cell.Value and isinstance(cell.Value, str):
                        cell.Value = translate_text(cell.Value, translation_dict, direction)
        
        new_output_path = output_path.replace('.xls', '.xlsx')
        wb.SaveAs(new_output_path, FileFormat=51)
        wb.Close()
        excel.Quit()
        pythoncom.CoUninitialize()
        
        if new_output_path != output_path:
            if os.path.exists(output_path):
                os.remove(output_path)
            os.rename(new_output_path, output_path)
        
        return True
    except Exception as e:
        logger.error(f".xls faylni tarjima qilishda xato: {e}")
        return False

def translate_old_ppt(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(input_path)
        
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if hasattr(shape, "TextFrame"):
                    if shape.TextFrame.HasText:
                        shape.TextFrame.TextRange.Text = translate_text(
                            shape.TextFrame.TextRange.Text, 
                            translation_dict, 
                            direction
                        )
        
        new_output_path = output_path.replace('.ppt', '.pptx')
        presentation.SaveAs(new_output_path, 24)
        presentation.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()
        
        if new_output_path != output_path:
            if os.path.exists(output_path):
                os.remove(output_path)
            os.rename(new_output_path, output_path)
        
        return True
    except Exception as e:
        logger.error(f".ppt faylni tarjima qilishda xato: {e}")
        return False

def translate_file(input_path: str, output_path: str, translation_dict: dict, direction: str, file_extension: str) -> bool:
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            if file_extension == '.doc':
                docx_path = os.path.join(temp_dir, "converted.docx")
                if not convert_doc_to_docx(input_path, docx_path):
                    return False
                return translate_docx(docx_path, output_path, translation_dict, direction)
            
            elif file_extension == '.xls':
                xlsx_path = os.path.join(temp_dir, "converted.xlsx")
                if not convert_xls_to_xlsx(input_path, xlsx_path):
                    return False
                return translate_excel(xlsx_path, output_path, translation_dict, direction)
            
            elif file_extension == '.ppt':
                pptx_path = os.path.join(temp_dir, "converted.pptx")
                if not convert_ppt_to_pptx(input_path, pptx_path):
                    return False
                return translate_pptx(pptx_path, output_path, translation_dict, direction)
            
            elif file_extension == '.docx':
                return translate_docx(input_path, output_path, translation_dict, direction)
            
            elif file_extension == '.xlsx':
                return translate_excel(input_path, output_path, translation_dict, direction)
            
            elif file_extension == '.pptx':
                return translate_pptx(input_path, output_path, translation_dict, direction)
            
            else:
                logger.error(f"Noto'g'ri fayl formati: {file_extension}")
                return False
                
    except Exception as e:
        logger.error(f"Fayl tarjimasida xatolik: {e}", exc_info=True)
        return False

def translate_docx(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        doc = docx.Document(input_path)
        
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text:
                    run.text = translate_text(run.text, translation_dict, direction)
        
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            if run.text:
                                run.text = translate_text(run.text, translation_dict, direction)
        
        doc.save(output_path)
        return True
    except Exception as e:
        logger.error(f"DOCX tarjimasida xato: {e}")
        return False

def translate_excel(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        wb = load_workbook(input_path)
        
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        cell.value = translate_text(cell.value, translation_dict, direction)
        
        wb.save(output_path)
        wb.close()
        return True
    except Exception as e:
        logger.error(f"Excel tarjimasida xato: {e}")
        return False

def translate_pptx(input_path: str, output_path: str, translation_dict: dict, direction: str) -> bool:
    try:
        prs = Presentation(input_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                run.text = translate_text(run.text, translation_dict, direction)
            
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if run.text:
                                            run.text = translate_text(run.text, translation_dict, direction)
        
        prs.save(output_path)
        return True
    except Exception as e:
        logger.error(f"PowerPoint tarjimasida xato: {e}")
        return False
    
def convert_doc_to_docx(input_path: str, output_path: str) -> bool:
    try:
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        doc = word.Documents.Open(input_path)
        doc.SaveAs2(output_path, FileFormat=16)
        doc.Close()
        word.Quit()
        pythoncom.CoUninitialize()
        return True
    except Exception as e:
        logger.error(f".doc -> .docx konvertatsiyada xato: {e}")
        return False

def convert_xls_to_xlsx(input_path: str, output_path: str) -> bool:
    try:
        pythoncom.CoInitialize()
        excel = win32com.client.Dispatch("Excel.Application")
        wb = excel.Workbooks.Open(input_path)
        wb.SaveAs(output_path, FileFormat=51)
        wb.Close()
        excel.Quit()
        pythoncom.CoUninitialize()
        return True
    except Exception as e:
        logger.error(f".xls -> .xlsx konvertatsiyada xato: {e}")
        return False

def convert_ppt_to_pptx(input_path: str, output_path: str) -> bool:
    try:
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(input_path)
        presentation.SaveAs(output_path, 24)
        presentation.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()
        return True
    except Exception as e:
        logger.error(f".ppt -> .pptx konvertatsiyada xato: {e}")
        return False
    
def translate_text(text: str, translation_dict: dict, direction: str) -> str:
    try:
        if not isinstance(text, str):
            return text
            
        if direction == 'lat_cyr':
            # Maxsus birikmalar uchun (katta harflar bilan boshlangan)
            text = re.sub(r'\bYe\b', '–ï', text)
            text = re.sub(r'Ye', '–ï', text)
            text = re.sub(r'\bYa\b', '–Ø', text)
            text = re.sub(r'Ya', '–Ø', text)
            text = re.sub(r'\bYu\b', '–Æ', text)
            text = re.sub(r'Yu', '–Æ', text)
            
            # Kichik harflar uchun
            text = re.sub(r'\bye\b', '–µ', text)
            text = re.sub(r'ye', '–µ', text)
            text = re.sub(r'\bya\b', '—è', text)
            text = re.sub(r'ya', '—è', text)
            text = re.sub(r'\byu\b', '—é', text)
            text = re.sub(r'yu', '—é', text)
            
            # Qolgan maxsus birikmalar (sh, ch, g', o')
            text = re.sub(r'\bSh\b', '–®', text)
            text = re.sub(r'\bSH\b', '–®', text)
            text = re.sub(r'\bSh', '–®', text)
            text = re.sub(r'SH', '–®', text)
            text = re.sub(r'sh', '—à', text)
            
            text = re.sub(r'\bCh\b', '–ß', text)
            text = re.sub(r'\bCH\b', '–ß', text)
            text = re.sub(r'\bCh', '–ß', text)
            text = re.sub(r'CH', '–ß', text)
            text = re.sub(r'ch', '—á', text)
            
            text = re.sub(r"g['‚Äò‚Äô`]", "“ì", text)
            text = re.sub(r"G['‚Äò‚Äô`]", "“í", text)
            text = re.sub(r"o['‚Äò‚Äô`]", "—û", text)
            text = re.sub(r"O['‚Äò‚Äô`]", "–é", text)
            text = re.sub(r"['‚Äò‚Äô`]", "—ä", text)
            
            # Qolgan harflar
            for lat, cyr in LATIN_TO_CYRILLIC.items():
                if len(lat) == 1 and lat not in ["'", "`"]:
                    text = text.replace(lat, cyr)
        else:
            # Kirill -> Lotin
            for cyr, lat in CYRILLIC_TO_LATIN.items():
                text = text.replace(cyr, lat)
        
        # Unicode matnni to'g'ri kodlash
        return text.encode('utf-8', 'ignore').decode('utf-8')
    
    except Exception as e:
        logger.error(f"Tarjima qilishda xato: {e}")
        return text
    
async def error_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    logger.error(msg="Exception while handling an update:", exc_info=context.error)
    if update and update.message:
        await update.message.reply_text(
            "‚ö†Ô∏è Kutilmagan xato yuz berdi. Iltimos, qayta urinib ko'ring yoki /start buyrug'ini yuboring."
        )

def main() -> None:
    application = Application.builder().token(TOKEN).build()
    
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CallbackQueryHandler(button_handler))
    
    docx_pdf_handler = MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_docx_pdf)
    translation_handler = MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_translation_doc)
    zip_handler = MessageHandler(
        (filters.Document.ALL | filters.PHOTO | filters.VIDEO | filters.AUDIO) & ~filters.COMMAND, 
        handle_zip_file
    )
    ocr_handler = MessageHandler(
        (filters.Document.ALL | filters.PHOTO) & ~filters.COMMAND, 
        handle_ocr
    )
    compress_handler = MessageHandler(
        (filters.Document.ALL | filters.PHOTO) & ~filters.COMMAND, 
        handle_compress
    )
    
    application.add_handler(docx_pdf_handler, group=1)
    application.add_handler(translation_handler, group=2)
    application.add_handler(zip_handler, group=3)
    application.add_handler(ocr_handler, group=4)
    application.add_handler(compress_handler, group=5)
    application.add_handler(MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_excel_pdf), group=6)
    application.add_handler(MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_ppt_pdf), group=7)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_qr_gen), group=8)
    application.add_handler(MessageHandler(filters.PHOTO & ~filters.COMMAND, handle_qr_scan), group=9)
    application.add_handler(MessageHandler(filters.Document.PDF & ~filters.COMMAND, handle_pdf_watermark), group=10)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_pdf_watermark_text), group=10)
    application.add_handler(MessageHandler(filters.Document.PDF & ~filters.COMMAND, handle_pdf_protect), group=11)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_pdf_protect_password), group=11)
    application.add_handler(MessageHandler(filters.Document.PDF & ~filters.COMMAND, handle_pdf_unprotect), group=12)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_pdf_unprotect_password), group=12)
    application.add_handler(MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_word_excel), group=13)

    
    application.add_handler(MessageHandler(filters.PHOTO & ~filters.COMMAND, handle_photo))
    application.add_error_handler(error_handler)
    
    application.run_polling()

if __name__ == "__main__":
    main()
