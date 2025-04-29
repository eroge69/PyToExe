import os
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext
from tkinter import ttk
from pptx import Presentation
import fitz  # PyMuPDF
from transformers import pipeline

# Inizializza pipeline di riassunto
summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# --- Funzioni di lettura file ---
def estrai_testo_pdf(percorso):
    testo = ""
    try:
        with fitz.open(percorso) as doc:
            for pagina in doc:
                testo += pagina.get_text()
    except Exception as e:
        print(f"Errore PDF {percorso}: {e}")
    return testo

def estrai_testo_ppt(percorso):
    testo = ""
    try:
        prs = Presentation(percorso)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    testo += shape.text + "\n"
    except Exception as e:
        print(f"Errore PPT {percorso}: {e}")
    return testo

def estrai_testo_txt(percorso):
    try:
        with open(percorso, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"Errore TXT {percorso}: {e}")
        return ""

def estrai_testo(percorso):
    if percorso.endswith(".pdf"):
        return estrai_testo_pdf(percorso)
    elif percorso.endswith(('.ppt', '.pptx')):
        return estrai_testo_ppt(percorso)
    elif percorso.endswith(".txt"):
        return estrai_testo_txt(percorso)
    return ""

# --- Funzione di ricerca ---
def cerca_keyword(cartelle, keyword):
    risultati = []
    for cartella in cartelle:
        for root, _, files in os.walk(cartella):
            for file in files:
                if file.endswith((".pdf", ".ppt", ".pptx", ".txt")):
                    percorso = os.path.join(root, file)
                    testo = estrai_testo(percorso)
                    if keyword.lower() in testo.lower():
                        try:
                            sunto = summarizer(testo[:1000], max_length=130, min_length=30, do_sample=False)[0]['summary_text']
                        except Exception as e:
                            sunto = f"Errore generazione sunto: {e}"
                        risultati.append((file, percorso, sunto))
    return risultati

# --- Interfaccia grafica ---
class RicercaApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Ricerca File con Sunto")

        self.cartelle = []

        frame_setup = ttk.LabelFrame(root, text="Setup cartelle")
        frame_setup.pack(fill="x", padx=10, pady=5)

        self.lista_cartelle = tk.Listbox(frame_setup, height=4)
        self.lista_cartelle.pack(side="left", fill="x", expand=True, padx=5, pady=5)

        btn_aggiungi = ttk.Button(frame_setup, text="Aggiungi cartella", command=self.aggiungi_cartella)
        btn_aggiungi.pack(side="right", padx=5, pady=5)

        frame_ricerca = ttk.LabelFrame(root, text="Ricerca")
        frame_ricerca.pack(fill="x", padx=10, pady=5)

        self.entry_keyword = ttk.Entry(frame_ricerca, width=50)
        self.entry_keyword.pack(side="left", padx=5, pady=5)

        btn_cerca = ttk.Button(frame_ricerca, text="Cerca", command=self.avvia_ricerca)
        btn_cerca.pack(side="left", padx=5, pady=5)

        self.area_risultati = scrolledtext.ScrolledText(root, height=20)
        self.area_risultati.pack(fill="both", expand=True, padx=10, pady=5)

    def aggiungi_cartella(self):
        cartella = filedialog.askdirectory()
        if cartella:
            self.cartelle.append(cartella)
            self.lista_cartelle.insert(tk.END, cartella)

    def avvia_ricerca(self):
        keyword = self.entry_keyword.get().strip()
        if not keyword:
            messagebox.showwarning("Attenzione", "Inserisci una parola chiave.")
            return

        self.area_risultati.delete(1.0, tk.END)
        risultati = cerca_keyword(self.cartelle, keyword)

        if not risultati:
            self.area_risultati.insert(tk.END, "Nessun risultato trovato.\n")
        else:
            for nome, percorso, sunto in risultati:
                self.area_risultati.insert(tk.END, f"File: {nome}\nPercorso: {percorso}\nSunto: {sunto}\n{'-'*60}\n")

# --- Avvio programma ---
if __name__ == "__main__":
    root = tk.Tk()
    app = RicercaApp(root)
    root.mainloop()
